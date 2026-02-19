#!/usr/bin/env python3
"""
MCP Video Tutorial Server
Provides video analysis tools, skill resources, and tutorial prompt templates
for the Tutorialinator workflow. Works as both a Claude Code MCP integration
and a standalone MCP server that other clients (e.g. Goose) can connect to.
"""

import os
import json
import subprocess
from pathlib import Path
from typing import Optional, Any

from mcp.server.fastmcp import FastMCP
from mcp.server.fastmcp.prompts.base import UserMessage, AssistantMessage
import whisper

# Initialize FastMCP server
mcp = FastMCP("tutorialinator")

# Skill directory: two levels up from this file → ~/.claude/skills/tutorialinator/
SKILL_DIR = Path(__file__).resolve().parent.parent.parent


def resolve_video_path(video_path: str) -> Path:
    """
    Resolve a video file path, handling macOS Unicode quirks.

    macOS screen recordings use U+202F (NARROW NO-BREAK SPACE) before AM/PM
    in filenames, but users/LLMs type regular spaces. This function tries
    the path as-is first, then attempts common substitutions.
    """
    p = Path(video_path)
    if p.exists():
        return p

    # Try replacing regular space with narrow no-break space (U+202F) and vice versa
    for find, replace in [(' ', '\u202f'), ('\u202f', ' ')]:
        alt = Path(video_path.replace(find, replace))
        if alt.exists():
            return alt

    # Try selective replacement: only around AM/PM patterns
    import re
    for pattern, repl in [
        (r' (AM|PM)', '\u202f\\1'),   # space before AM/PM → NNBSP
        (r'\u202f(AM|PM)', ' \\1'),   # NNBSP before AM/PM → space
    ]:
        alt_str = re.sub(pattern, repl, video_path)
        alt = Path(alt_str)
        if alt.exists():
            return alt

    # Return original path (will fail with "not found" downstream)
    return p

# Configuration
CACHE_DIR = Path(os.getenv("VIDEO_TUTORIAL_CACHE", Path.home() / ".cache" / "video-tutorial"))
CACHE_DIR.mkdir(parents=True, exist_ok=True)

WHISPER_MODEL_SIZE = os.getenv("WHISPER_MODEL_SIZE", "base")
WHISPER_DEVICE = os.getenv("WHISPER_DEVICE", "cpu")

# Load Whisper model lazily on first use
whisper_model = None

def get_whisper_model():
    global whisper_model
    if whisper_model is None:
        try:
            print(f"Loading Whisper model: {WHISPER_MODEL_SIZE}")
            whisper_model = whisper.load_model(WHISPER_MODEL_SIZE, device=WHISPER_DEVICE)
            print("Whisper model loaded successfully")
        except Exception as e:
            print(f"Warning: Failed to load Whisper model: {e}")
            print("Transcription will not be available. Run auto-setup.sh to fix.")
            raise
    return whisper_model


@mcp.tool()
def download_video(url: str, quality: str = "1080p") -> dict[str, Any]:
    """
    Download video from URL (YouTube, Vimeo, etc.) using yt-dlp.

    Args:
        url: Video URL from supported platforms
        quality: Desired quality (1080p, 720p, 480p, etc.)

    Returns:
        Dictionary with local_path and metadata
    """
    try:
        output_dir = CACHE_DIR / "videos"
        output_dir.mkdir(exist_ok=True)

        # Generate output filename based on video ID
        output_template = str(output_dir / "%(id)s.%(ext)s")

        # Use yt-dlp to download video
        cmd = [
            "yt-dlp",
            "-f", f"bestvideo[height<={quality[:-1]}]+bestaudio/best",
            "-o", output_template,
            "--write-info-json",
            url
        ]

        result = subprocess.run(cmd, capture_output=True, text=True, check=True)

        # Find the downloaded file
        video_files = list(output_dir.glob("*"))
        video_file = max(video_files, key=lambda x: x.stat().st_mtime)

        # Load metadata from info.json if available
        info_json = video_file.with_suffix(".info.json")
        metadata = {}
        if info_json.exists():
            with open(info_json) as f:
                info = json.load(f)
                metadata = {
                    "duration": info.get("duration"),
                    "resolution": f"{info.get('width')}x{info.get('height')}",
                    "fps": info.get("fps"),
                    "title": info.get("title"),
                    "uploader": info.get("uploader"),
                }

        return {
            "success": True,
            "local_path": str(video_file),
            "metadata": metadata
        }

    except subprocess.CalledProcessError as e:
        return {
            "success": False,
            "error": f"Download failed: {e.stderr}"
        }
    except Exception as e:
        return {
            "success": False,
            "error": str(e)
        }


@mcp.tool()
def transcribe_with_timestamps(
    video_path: str,
    language: str = "auto"
) -> dict[str, Any]:
    """
    Transcribe video with word-level timestamps using Whisper.

    Args:
        video_path: Path to video file
        language: Language code or 'auto' for detection

    Returns:
        Structured transcript with segments and word-level timestamps
    """
    try:
        # Check if enhanced whisper-timestamped is available
        try:
            import whisper_timestamped as whisper_ts
            use_enhanced = True
        except ImportError:
            use_enhanced = False
            print("whisper-timestamped not available, using standard Whisper (segment-level only)")

        video_path = resolve_video_path(video_path)
        if not video_path.exists():
            return {"success": False, "error": f"Video file not found: {video_path}"}

        # Load audio
        if use_enhanced:
            audio = whisper_ts.load_audio(str(video_path))
            result = whisper_ts.transcribe(
                get_whisper_model(),
                audio,
                language=None if language == "auto" else language
            )

            # Extract word-level timestamps
            words = []
            for segment in result["segments"]:
                for word in segment.get("words", []):
                    words.append({
                        "word": word["text"],
                        "start_time": word["start"],
                        "end_time": word["end"],
                        "confidence": word.get("confidence", 1.0)
                    })
        else:
            # Fallback to standard Whisper (segment-level only)
            result = get_whisper_model().transcribe(
                str(video_path),
                language=None if language == "auto" else language
            )
            words = []  # No word-level timestamps available

        segments = [
            {
                "start_time": seg["start"],
                "end_time": seg["end"],
                "text": seg["text"].strip(),
            }
            for seg in result["segments"]
        ]

        full_text = " ".join(seg["text"].strip() for seg in result["segments"])

        return {
            "success": True,
            "segments": segments,
            "words": words,
            "full_text": full_text,
            "language": result.get("language", "unknown"),
            "has_word_timestamps": use_enhanced
        }

    except Exception as e:
        return {
            "success": False,
            "error": f"Transcription failed: {str(e)}"
        }


@mcp.tool()
def detect_scenes(
    video_path: str,
    threshold: float = 0.3,
    min_scene_duration: float = 2.0
) -> dict[str, Any]:
    """
    Detect scene changes in video using FFmpeg.

    Args:
        video_path: Path to video file
        threshold: Scene change sensitivity (0.0-1.0, lower=more sensitive)
        min_scene_duration: Minimum scene length in seconds

    Returns:
        List of scene boundaries with timestamps
    """
    try:
        video_path = resolve_video_path(video_path)
        if not video_path.exists():
            return {"success": False, "error": f"Video file not found: {video_path}"}

        # Use FFmpeg scene detection filter
        cmd = [
            "ffmpeg",
            "-i", str(video_path),
            "-filter:v", f"select='gt(scene,{threshold})',showinfo",
            "-f", "null",
            "-"
        ]

        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)

        # Parse scene timestamps from ffmpeg output
        scenes = []
        last_time = 0.0

        for line in result.stdout.split('\n'):
            if 'pts_time:' in line:
                # Extract timestamp
                parts = line.split('pts_time:')
                if len(parts) > 1:
                    time_str = parts[1].split()[0]
                    try:
                        time = float(time_str)
                        if time - last_time >= min_scene_duration:
                            scenes.append({
                                "start_time": last_time,
                                "end_time": time,
                                "duration": time - last_time
                            })
                            last_time = time
                    except ValueError:
                        continue

        # If no scenes detected, return whole video as one scene
        if not scenes:
            # Get video duration
            probe_cmd = [
                "ffprobe",
                "-v", "error",
                "-show_entries", "format=duration",
                "-of", "default=noprint_wrappers=1:nokey=1",
                str(video_path)
            ]
            duration_result = subprocess.run(probe_cmd, capture_output=True, text=True)
            try:
                duration = float(duration_result.stdout.strip())
                scenes = [{
                    "start_time": 0.0,
                    "end_time": duration,
                    "duration": duration
                }]
            except ValueError:
                return {"success": False, "error": "Could not determine video duration"}

        return {
            "success": True,
            "scenes": scenes,
            "scene_count": len(scenes)
        }

    except Exception as e:
        return {
            "success": False,
            "error": f"Scene detection failed: {str(e)}"
        }


@mcp.tool()
def extract_key_frames(
    video_path: str,
    timestamps: Optional[list[float]] = None,
    scene_data: Optional[dict] = None
) -> dict[str, Any]:
    """
    Extract key frames from video at specified timestamps or scene boundaries.

    Args:
        video_path: Path to video file
        timestamps: List of timestamps in seconds (optional)
        scene_data: Output from detect_scenes (optional, uses scene boundaries)

    Returns:
        List of extracted frame paths with timestamps
    """
    try:
        video_path = resolve_video_path(video_path)
        if not video_path.exists():
            return {"success": False, "error": f"Video file not found: {video_path}"}

        # Determine timestamps to extract
        extract_times = []
        if timestamps:
            extract_times = timestamps
        elif scene_data and scene_data.get("success"):
            # Extract frames at scene boundaries
            extract_times = [scene["start_time"] for scene in scene_data["scenes"]]
        else:
            return {"success": False, "error": "No timestamps or scene data provided"}

        # Create output directory
        output_dir = CACHE_DIR / "frames" / video_path.stem
        output_dir.mkdir(parents=True, exist_ok=True)

        frames = []
        for i, timestamp in enumerate(extract_times):
            output_file = output_dir / f"frame_{i:04d}_{timestamp:.2f}s.jpg"

            cmd = [
                "ffmpeg",
                "-ss", str(timestamp),
                "-i", str(video_path),
                "-vframes", "1",
                "-q:v", "2",
                "-y",
                str(output_file)
            ]

            subprocess.run(cmd, capture_output=True, check=True)

            frames.append({
                "timestamp": timestamp,
                "frame_path": str(output_file),
                "frame_number": i
            })

        return {
            "success": True,
            "frames": frames,
            "frame_count": len(frames)
        }

    except subprocess.CalledProcessError as e:
        return {
            "success": False,
            "error": f"Frame extraction failed: {e.stderr.decode() if e.stderr else str(e)}"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"Frame extraction failed: {str(e)}"
        }


@mcp.tool()
def detect_code_in_frames(
    frame_paths: list[str],
    languages: Optional[list[str]] = None
) -> dict[str, Any]:
    """
    Detect code snippets in video frames using OCR.

    Args:
        frame_paths: List of paths to frame images
        languages: Programming languages to detect (optional)

    Returns:
        List of detected code snippets with metadata
    """
    try:
        # Check if RapidOCR is available
        try:
            from rapidocr_onnxruntime import RapidOCR
            ocr = RapidOCR()
            ocr_available = True
        except ImportError:
            ocr_available = False
            return {
                "success": False,
                "error": "RapidOCR not installed. Install with: pip install rapidocr-onnxruntime",
                "recommendation": "OCR capabilities require optional dependencies. See setup documentation."
            }

        code_snippets = []

        for frame_path in frame_paths:
            if not Path(frame_path).exists():
                continue

            # Run OCR on frame
            result, elapse = ocr(frame_path)

            if result:
                # Extract text from OCR results
                text_blocks = [item[1] for item in result]
                full_text = "\n".join(text_blocks)

                # Simple heuristic: code usually has { } ; indentation
                code_indicators = ['{', '}', ';', '  ', 'function', 'const', 'let', 'var', 'import', 'export']
                code_score = sum(1 for indicator in code_indicators if indicator in full_text)

                if code_score >= 3:  # Threshold for "looks like code"
                    code_snippets.append({
                        "frame_path": frame_path,
                        "code_text": full_text,
                        "confidence": code_score / len(code_indicators),
                        "detected_keywords": [ind for ind in code_indicators if ind in full_text]
                    })

        return {
            "success": True,
            "code_snippets": code_snippets,
            "snippet_count": len(code_snippets)
        }

    except Exception as e:
        return {
            "success": False,
            "error": f"Code detection failed: {str(e)}"
        }


@mcp.tool()
def generate_chapters(
    transcript: dict,
    scenes: dict,
    min_chapter_length: float = 60.0
) -> dict[str, Any]:
    """
    Generate intelligent chapter suggestions from transcript and scene data.

    Args:
        transcript: Output from transcribe_with_timestamps
        scenes: Output from detect_scenes
        min_chapter_length: Minimum chapter duration in seconds

    Returns:
        Suggested chapter structure with titles and summaries
    """
    try:
        if not transcript.get("success") or not scenes.get("success"):
            return {"success": False, "error": "Invalid transcript or scene data"}

        segments = transcript["segments"]
        scene_list = scenes["scenes"]

        chapters = []
        current_chapter_start = 0.0
        current_chapter_text = []

        for scene in scene_list:
            scene_start = scene["start_time"]
            scene_end = scene["end_time"]

            # Check if enough time has passed for a new chapter
            if scene_start - current_chapter_start >= min_chapter_length:
                # Find segments in this chapter
                chapter_segments = [
                    seg for seg in segments
                    if seg["start_time"] >= current_chapter_start and seg["end_time"] <= scene_start
                ]

                if chapter_segments:
                    chapter_text = " ".join(seg["text"] for seg in chapter_segments)

                    # Generate chapter title from first sentence
                    first_sentence = chapter_text.split('.')[0][:100]
                    title = first_sentence if first_sentence else f"Chapter {len(chapters) + 1}"

                    chapters.append({
                        "start_time": current_chapter_start,
                        "end_time": scene_start,
                        "title": title.strip(),
                        "summary": chapter_text[:200] + "..." if len(chapter_text) > 200 else chapter_text,
                        "duration": scene_start - current_chapter_start
                    })

                    current_chapter_start = scene_start

        # Add final chapter
        final_segments = [
            seg for seg in segments
            if seg["start_time"] >= current_chapter_start
        ]
        if final_segments:
            chapter_text = " ".join(seg["text"] for seg in final_segments)
            chapters.append({
                "start_time": current_chapter_start,
                "end_time": segments[-1]["end_time"],
                "title": chapter_text.split('.')[0][:100].strip() or f"Chapter {len(chapters) + 1}",
                "summary": chapter_text[:200] + "..." if len(chapter_text) > 200 else chapter_text,
                "duration": segments[-1]["end_time"] - current_chapter_start
            })

        return {
            "success": True,
            "chapters": chapters,
            "chapter_count": len(chapters)
        }

    except Exception as e:
        return {
            "success": False,
            "error": f"Chapter generation failed: {str(e)}"
        }


@mcp.tool()
def extract_slides(file_path: str) -> dict[str, Any]:
    """
    Extract slide content from presentation files (.pptx or .pdf).

    Args:
        file_path: Path to a .pptx or .pdf file

    Returns:
        List of slide dicts with slide_number, title, body, speaker_notes, and image_b64 (PDF only)
    """
    try:
        p = Path(file_path)
        if not p.exists():
            return {"success": False, "error": f"File not found: {file_path}"}

        suffix = p.suffix.lower()

        if suffix == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return {
                    "success": False,
                    "error": "python-pptx not installed. Install with: pip install python-pptx",
                }

            prs = Presentation(str(p))
            slides = []
            for i, slide in enumerate(prs.slides, start=1):
                title = ""
                body_parts = []
                title_shape_id = slide.shapes.title.shape_id if slide.shapes.title else None
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        if title_shape_id is not None and shape.shape_id == title_shape_id:
                            title = shape.text_frame.text
                        else:
                            body_parts.append(shape.text_frame.text)

                # Fallback: if title still empty, try the title placeholder directly
                if not title and slide.shapes.title:
                    title = slide.shapes.title.text

                notes = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text

                slides.append({
                    "slide_number": i,
                    "title": title,
                    "body": "\n".join(body_parts),
                    "speaker_notes": notes,
                    "image_b64": None,
                })

            return {"success": True, "slides": slides, "slide_count": len(slides)}

        elif suffix == ".pdf":
            try:
                import fitz  # PyMuPDF
            except ImportError:
                return {
                    "success": False,
                    "error": "PyMuPDF not installed. Install with: pip install PyMuPDF",
                }

            import base64

            doc = fitz.open(str(p))
            slides = []
            for i, page in enumerate(doc, start=1):
                text = page.get_text().strip()

                # Render page to JPEG
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("jpeg")
                img_b64 = base64.b64encode(img_bytes).decode("ascii")

                # Use first line as title heuristic
                lines = text.split("\n") if text else []
                title = lines[0] if lines else ""
                body = "\n".join(lines[1:]) if len(lines) > 1 else ""

                slides.append({
                    "slide_number": i,
                    "title": title,
                    "body": body,
                    "speaker_notes": "",
                    "image_b64": img_b64,
                })
            doc.close()

            return {"success": True, "slides": slides, "slide_count": len(slides)}

        else:
            return {
                "success": False,
                "error": f"Unsupported file format: {suffix}. Supported: .pptx, .pdf",
            }

    except Exception as e:
        return {
            "success": False,
            "error": f"Slide extraction failed: {str(e)}",
        }


@mcp.tool()
def get_video_metadata(video_path: str) -> dict[str, Any]:
    """
    Extract comprehensive metadata from video file.

    Args:
        video_path: Path to video file

    Returns:
        Video metadata including codec, resolution, duration, etc.
    """
    try:
        video_path = resolve_video_path(video_path)
        if not video_path.exists():
            return {"success": False, "error": f"Video file not found: {video_path}"}

        cmd = [
            "ffprobe",
            "-v", "quiet",
            "-print_format", "json",
            "-show_format",
            "-show_streams",
            str(video_path)
        ]

        result = subprocess.run(cmd, capture_output=True, text=True, check=True)
        data = json.loads(result.stdout)

        # Extract relevant metadata
        video_stream = next((s for s in data["streams"] if s["codec_type"] == "video"), None)
        audio_stream = next((s for s in data["streams"] if s["codec_type"] == "audio"), None)
        format_data = data.get("format", {})

        metadata = {
            "duration": float(format_data.get("duration", 0)),
            "size_bytes": int(format_data.get("size", 0)),
            "bitrate": int(format_data.get("bit_rate", 0)),
        }

        if video_stream:
            metadata.update({
                "video_codec": video_stream.get("codec_name"),
                "width": video_stream.get("width"),
                "height": video_stream.get("height"),
                "fps": (lambda r: int(r.split("/")[0]) / int(r.split("/")[1]) if "/" in r and int(r.split("/")[1]) != 0 else float(r) if r else 0)(video_stream.get("r_frame_rate", "0")),
                "resolution": f"{video_stream.get('width')}x{video_stream.get('height')}"
            })

        if audio_stream:
            metadata.update({
                "audio_codec": audio_stream.get("codec_name"),
                "sample_rate": audio_stream.get("sample_rate"),
                "channels": audio_stream.get("channels")
            })

        return {
            "success": True,
            "metadata": metadata
        }

    except subprocess.CalledProcessError as e:
        return {
            "success": False,
            "error": f"Metadata extraction failed: {e.stderr}"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"Metadata extraction failed: {str(e)}"
        }


# ---------------------------------------------------------------------------
# MCP Resources — expose skill files so any MCP client can read them
# ---------------------------------------------------------------------------

def _read_skill_file(relative_path: str) -> str:
    """Read a file from the skill directory, with a clear error if missing."""
    path = SKILL_DIR / relative_path
    if not path.exists():
        return f"Error: {relative_path} not found at {path}"
    return path.read_text()


@mcp.resource(
    "tutorialinator://skill",
    name="skill",
    description="Full Tutorialinator design system, pedagogy spec, widget patterns, and orchestration rules (~1275 lines)",
    mime_type="text/markdown",
)
def resource_skill() -> str:
    return _read_skill_file("SKILL.md")


@mcp.resource(
    "tutorialinator://design-updates",
    name="design-updates",
    description="V3 design system changelog — what changed from the previous architecture",
    mime_type="text/markdown",
)
def resource_design_updates() -> str:
    return _read_skill_file("templates/DESIGN_UPDATES.md")


@mcp.resource(
    "tutorialinator://templates-guide",
    name="templates-guide",
    description="Overview of the single-HTML-file template architecture",
    mime_type="text/markdown",
)
def resource_templates_guide() -> str:
    return _read_skill_file("templates/README.md")


# ---------------------------------------------------------------------------
# MCP Prompts — ready-to-use prompt templates for each tutorial mode
# ---------------------------------------------------------------------------

@mcp.prompt(
    name="tutorial_topic",
    description="Generate an interactive tutorial on any topic using the Tutorialinator design system (Topic Mode)",
)
def prompt_tutorial_topic(topic: str) -> list:
    skill = _read_skill_file("SKILL.md")
    return [
        AssistantMessage(
            "I have the Tutorialinator design system loaded. "
            "I will follow it exactly to produce a single self-contained HTML tutorial.\n\n"
            + skill
        ),
        UserMessage(
            f"Using the Tutorialinator design system above, generate an interactive tutorial about: {topic}\n\n"
            "Follow the Topic Mode Workflow from the skill specification."
        ),
    ]


@mcp.prompt(
    name="tutorial_video",
    description="Generate an interactive tutorial from a local video file (Video Mode)",
)
def prompt_tutorial_video(video_path: str) -> list:
    skill = _read_skill_file("SKILL.md")
    return [
        AssistantMessage(
            "I have the Tutorialinator design system loaded. "
            "I will follow it exactly to produce a single self-contained HTML tutorial.\n\n"
            + skill
        ),
        UserMessage(
            f"Using the Tutorialinator design system above, generate an interactive tutorial from this video: {video_path}\n\n"
            "Follow the Video Mode Workflow from the skill specification. "
            "Use the MCP video tools (transcribe_with_timestamps, detect_scenes, extract_key_frames, etc.) to analyze the video first."
        ),
    ]


@mcp.prompt(
    name="tutorial_research",
    description="Generate an interactive tutorial synthesized from provided URLs (Research Mode)",
)
def prompt_tutorial_research(urls: str) -> list:
    skill = _read_skill_file("SKILL.md")
    return [
        AssistantMessage(
            "I have the Tutorialinator design system loaded. "
            "I will follow it exactly to produce a single self-contained HTML tutorial.\n\n"
            + skill
        ),
        UserMessage(
            f"Using the Tutorialinator design system above, generate an interactive tutorial from these resources:\n{urls}\n\n"
            "Follow the Research Mode Workflow from the skill specification. "
            "Fetch each URL, synthesize the content, and build a tutorial with inline citations."
        ),
    ]


@mcp.prompt(
    name="tutorial_deep_research",
    description="Generate a comprehensive interactive tutorial via web research (Deep Research Mode)",
)
def prompt_tutorial_deep_research(topic: str) -> list:
    skill = _read_skill_file("SKILL.md")
    return [
        AssistantMessage(
            "I have the Tutorialinator design system loaded. "
            "I will follow it exactly to produce a single self-contained HTML tutorial.\n\n"
            + skill
        ),
        UserMessage(
            f"Using the Tutorialinator design system above, do deep research and generate a comprehensive interactive tutorial about: {topic}\n\n"
            "Follow the Deep Research Mode Workflow from the skill specification. "
            "Search the web for 5-15 quality sources, build a concept map, identify misconceptions, then generate the tutorial."
        ),
    ]


# ---------------------------------------------------------------------------
# Server entry point
# ---------------------------------------------------------------------------

def run_server():
    """Entry point for running the MCP server."""
    print("Starting MCP Video Tutorial Server...")
    print(f"Cache directory: {CACHE_DIR}")
    print(f"Whisper model: {WHISPER_MODEL_SIZE} on {WHISPER_DEVICE}")
    print(f"Skill directory: {SKILL_DIR}")
    mcp.run()


if __name__ == "__main__":
    run_server()
